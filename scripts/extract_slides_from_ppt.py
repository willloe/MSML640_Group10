import argparse
import logging
from pathlib import Path
from typing import List
import os

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("python-pptx not installed.")

try:
    from PIL import Image
    import io
    HAS_PIL = True
except ImportError:
    HAS_PIL = False
    print("Pillow not installed.")

try:
    from pdf2image import convert_from_path
    import tempfile
    HAS_PDF2IMAGE = True
except ImportError:
    HAS_PDF2IMAGE = False

try:
    import win32com.client
    import pythoncom
    HAS_WIN32COM = True
except ImportError:
    HAS_WIN32COM = False

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class PowerPointExtractor:

    def __init__(self, output_dir: Path, method: str = 'auto'):

        self.output_dir = output_dir
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.method = method

        if method == 'auto':
            if HAS_WIN32COM:
                self.method = 'win32com'
                logger.info("Using win32com method")
            elif HAS_PDF2IMAGE:
                self.method = 'pdf'
                logger.info("Using PDF conversion method")
            else:
                logger.warning("No optimal extraction method available.")
                self.method = 'basic'

    def extract_with_win32com(self, ppt_path: Path, prefix: str = None) -> List[Path]:

        if not HAS_WIN32COM:
            raise RuntimeError("win32com not available.")

        pythoncom.CoInitialize()

        try:
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            powerpoint.Visible = 1

            presentation = powerpoint.Presentations.Open(str(ppt_path.absolute()), WithWindow=False)

            output_files = []
            prefix = prefix or ppt_path.stem

            for i, slide in enumerate(presentation.Slides):
                output_file = self.output_dir / f"{prefix}_slide_{i+1:03d}.png"

                slide.Export(str(output_file.absolute()), "PNG", 1024, 768)

                output_files.append(output_file)
                logger.info(f"Extracted slide {i+1}/{presentation.Slides.Count}")

            presentation.Close()
            powerpoint.Quit()

            return output_files

        finally:
            pythoncom.CoUninitialize()

    def extract_with_pdf(self, ppt_path: Path, prefix: str = None) -> List[Path]:

        if not HAS_PDF2IMAGE:
            raise RuntimeError("pdf2image not available.")

        if not HAS_WIN32COM and not HAS_PPTX:
            raise RuntimeError("Need either win32com or python-pptx to convert to PDF")

        prefix = prefix or ppt_path.stem


        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp_pdf:
            tmp_pdf_path = Path(tmp_pdf.name)

        try:

            if HAS_WIN32COM:
                pythoncom.CoInitialize()
                try:
                    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                    presentation = powerpoint.Presentations.Open(str(ppt_path.absolute()), WithWindow=False)
                    presentation.SaveAs(str(tmp_pdf_path.absolute()), 32)  # 32 = ppSaveAsPDF
                    presentation.Close()
                    powerpoint.Quit()
                finally:
                    pythoncom.CoUninitialize()
            else:
                logger.warning("Cannot convert to PDF without win32com.")
                return self.extract_basic(ppt_path, prefix)

            images = convert_from_path(tmp_pdf_path, dpi=150, fmt='png')

            output_files = []
            for i, image in enumerate(images):
                output_file = self.output_dir / f"{prefix}_slide_{i+1:03d}.png"

                image = image.resize((1024, 768), Image.Resampling.LANCZOS)
                image.save(output_file, 'PNG')

                output_files.append(output_file)
                logger.info(f"Extracted slide {i+1}/{len(images)}")

            return output_files

        finally:
            if tmp_pdf_path.exists():
                tmp_pdf_path.unlink()

    def extract_basic(self, ppt_path: Path, prefix: str = None) -> List[Path]:
        if not HAS_PPTX:
            raise RuntimeError("python-pptx not available.")

        logger.warning("Basic method only extracts embedded images")

        presentation = Presentation(str(ppt_path))
        output_files = []
        prefix = prefix or ppt_path.stem

        for slide_idx, slide in enumerate(presentation.slides):
            image_count = 0
            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    image = shape.image
                    image_bytes = image.blob

                    output_file = self.output_dir / f"{prefix}_slide_{slide_idx+1:03d}_img_{image_count:02d}.png"

                    with open(output_file, 'wb') as f:
                        f.write(image_bytes)

                    output_files.append(output_file)
                    image_count += 1

            if image_count > 0:
                logger.info(f"Extracted {image_count} images from slide {slide_idx+1}")

        if not output_files:
            logger.warning("No images found in presentation.")

        return output_files

    def extract_presentation(self, ppt_path: Path, prefix: str = None) -> List[Path]:

        logger.info(f"Extracting slides from {ppt_path.name} using {self.method} method")

        if self.method == 'win32com':
            return self.extract_with_win32com(ppt_path, prefix)
        elif self.method == 'pdf':
            return self.extract_with_pdf(ppt_path, prefix)
        else:
            return self.extract_basic(ppt_path, prefix)


def main():
    parser = argparse.ArgumentParser(description='Extract slide images from PowerPoint files')
    parser.add_argument('--input', type=str, required=True,
                        help='Input PowerPoint file or directory')
    parser.add_argument('--output', type=str, default='data/raw',
                        help='Output directory for extracted images')
    parser.add_argument('--method', type=str, default='auto',
                        choices=['auto', 'win32com', 'pdf', 'basic'],
                        help='Extraction method to use')
    parser.add_argument('--prefix', type=str, default=None,
                        help='Prefix for output files (default: use filename)')

    args = parser.parse_args()

    input_path = Path(args.input)
    output_dir = Path(args.output)

    if not input_path.exists():
        logger.error(f"Input path not found: {input_path}")
        return

    extractor = PowerPointExtractor(output_dir, method=args.method)

    if input_path.is_file():
        ppt_files = [input_path]
    else:
        ppt_files = list(input_path.glob('*.pptx')) + list(input_path.glob('*.ppt'))

    if not ppt_files:
        logger.error(f"No PowerPoint files found in {input_path}")
        return

    logger.info(f"Found {len(ppt_files)} PowerPoint file(s)")

    all_outputs = []
    for ppt_file in ppt_files:
        try:
            prefix = args.prefix or ppt_file.stem
            outputs = extractor.extract_presentation(ppt_file, prefix)
            all_outputs.extend(outputs)
        except Exception as e:
            logger.error(f"Failed to extract from {ppt_file.name}: {e}")
            import traceback
            traceback.print_exc()
            continue

    logger.info(f"\nExtracted {len(all_outputs)} slide images to {output_dir}")
    logger.info(f"Extracted slides saved in: {output_dir}")


if __name__ == '__main__':
    main()
